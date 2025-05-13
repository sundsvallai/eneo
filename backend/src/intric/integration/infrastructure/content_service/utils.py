"""
Utility functions for processing different file types in integration services.
"""

import io
import re
import zipfile
from logging import getLogger
from typing import Dict, Optional, Tuple
import pandas as pd

logger = getLogger(__name__)


def binary_to_text(binary_data: bytes) -> str:
    """
    Extract readable text from binary data, properly handling Office documents.
    For Office documents (docx, xlsx, pptx), attempts to identify and extract
    the text content without the binary parts.

    Args:
        binary_data: The binary data to extract text from

    Returns:
        String containing the readable text portions
    """
    # Check if this looks like a ZIP file (Office documents are ZIP archives)
    if binary_data.startswith(b"PK\x03\x04"):
        try:
            import io
            import re
            import zipfile

            # Office documents are ZIP files with XML content
            with zipfile.ZipFile(io.BytesIO(binary_data)) as zip_file:
                # Check which type of Office document this is
                if "word/document.xml" in zip_file.namelist():
                    return extract_text_from_docx(binary_data, zip_file)
                elif any(
                    name.startswith("ppt/slides/slide") for name in zip_file.namelist()
                ):
                    return extract_text_from_pptx(binary_data, zip_file)

                # Generic XML extraction as a fallback for unknown Office formats
                all_text = []
                for filename in zip_file.namelist():
                    if filename.endswith(".xml"):
                        try:
                            with zip_file.open(filename) as xml_file:
                                content = xml_file.read()
                                # Extract text between XML tags
                                text_content = re.findall(b">([^<]+)</", content)
                                for text in text_content:
                                    try:
                                        decoded = text.decode("utf-8", errors="replace")
                                        if (
                                            len(decoded.strip()) > 5
                                        ):  # Only keep if not too short
                                            all_text.append(decoded.strip())
                                    except Exception:
                                        pass
                        except Exception:
                            continue

                if all_text:
                    return "\n".join(all_text)

        except Exception as e:
            logger.warning(f"Failed to extract text from Office document: {e}")
            # Fall back to the standard method if ZIP extraction fails

    # Standard method for non-Office documents or if Office parsing failed
    result = []
    current_text = ""

    # Filter out null bytes (cause PostgreSQL UTF-8 encoding errors)
    filtered_data = bytes(b for b in binary_data if b != 0)

    try:
        # Try to decode the filtered data as UTF-8
        decoded = filtered_data.decode("utf-8", errors="replace")

        # Process decoded text character by character
        for char in decoded:
            if char.isprintable() or char.isspace():
                current_text += char
            else:
                if len(current_text) > 20:
                    result.append(current_text.strip())
                current_text = ""

        # Add the last text chunk if substantial
        if len(current_text) > 20:
            result.append(current_text.strip())

        return "\n\n".join(result) if result else "[No readable text found]"

    except Exception:
        # Fallback to ASCII-only extraction
        for byte in filtered_data:
            if 32 <= byte <= 126:  # ASCII printable range
                current_text += chr(byte)
            else:
                if len(current_text) > 20:
                    result.append(current_text)
                current_text = ""

        if current_text and len(current_text) > 20:
            result.append(current_text)

        return "\n".join(result) if result else "[No readable text found]"


def extract_text_from_docx(
    binary_data: bytes, zip_file: Optional[zipfile.ZipFile] = None
) -> str:
    """
    Extract text from a Word document (DOCX format).

    Args:
        binary_data: Binary data of the DOCX file
        zip_file: Optional already opened ZipFile object

    Returns:
        Extracted text content
    """
    # If zip_file wasn't provided, create one
    if zip_file is None:
        zip_file = zipfile.ZipFile(io.BytesIO(binary_data))
        close_zip = True
    else:
        close_zip = False

    try:
        with zip_file.open("word/document.xml") as docx_file:
            xml_content = docx_file.read()

            # Look for text between <w:t> tags
            text_matches = re.findall(b"<w:t[^>]*>(.*?)</w:t>", xml_content, re.DOTALL)
            if text_matches:
                clean_text = []
                for match in text_matches:
                    try:
                        decoded = match.decode("utf-8", errors="replace")
                        if decoded.strip():
                            clean_text.append(decoded)
                    except Exception:
                        pass

                if clean_text:
                    return " ".join(clean_text)
    except Exception as e:
        logger.warning(f"Error extracting DOCX content: {e}")

    finally:
        if close_zip:
            zip_file.close()

    # Empty content
    return ""


def extract_text_from_pptx(
    binary_data: bytes, zip_file: Optional[zipfile.ZipFile] = None
) -> str:
    """
    Extract text from a PowerPoint presentation (PPTX format).

    Args:
        binary_data: Binary data of the PPTX file
        zip_file: Optional already opened ZipFile object

    Returns:
        Extracted text content
    """
    # If zip_file wasn't provided, create one
    if zip_file is None:
        zip_file = zipfile.ZipFile(io.BytesIO(binary_data))
        close_zip = True
    else:
        close_zip = False

    try:
        # Try using python-pptx if available
        try:
            from pptx import Presentation

            # Create a temporary file for the presentation
            with io.BytesIO(binary_data) as temp_pptx:
                # Load the presentation
                presentation = Presentation(temp_pptx)

                # Extract text from each slide
                slide_texts = []

                for slide_num, slide in enumerate(presentation.slides, 1):
                    # Extract text from all shapes on the slide
                    slide_content = []

                    # Process each shape
                    for shape in slide.shapes:
                        # Extract text from text frames
                        if hasattr(shape, "text") and shape.text.strip():
                            # Skip template placeholders
                            text = shape.text.strip()

                            # Skip common placeholders and formatting text
                            skip_phrases = [
                                "Click to edit",
                                "Double click",
                                "Click icon",
                                "Master",
                                "Title style",
                                "level",
                                "office theme",
                                "Microsoft Office",
                            ]

                            if any(phrase in text for phrase in skip_phrases):
                                continue

                            # Add the text with context
                            slide_content.append(text)

                    # Add the slide content if it's not empty
                    if slide_content:
                        # Mark the slide number
                        slide_text = f"Slide {slide_num}:\n" + "\n".join(slide_content)
                        slide_texts.append(slide_text)

                # Join all slide texts with clear separation
                if slide_texts:
                    return "\n\n".join(slide_texts)

        except ImportError:
            logger.warning(
                "python-pptx not available, falling back to basic extraction"
            )
        except Exception as e:
            logger.warning(f"Error extracting PPTX with python-pptx: {e}")

        # If python-pptx extraction failed, try a simpler approach
        slide_texts = []
        slide_files = sorted(
            [
                name
                for name in zip_file.namelist()
                if name.startswith("ppt/slides/slide")
            ]
        )

        for i, slide_file in enumerate(slide_files, 1):
            try:
                with zip_file.open(slide_file) as f:
                    slide_content = f.read()

                    # Extract paragraph text using more precise targeting
                    # This focuses on actual content and avoids placeholders
                    paragraphs = []

                    # Extract actual slide content from more specific XML paths
                    # 1. Extract text from title placeholders
                    title_matches = re.findall(
                        b"<p:title[^>]*>.*?<a:t[^>]*>(.*?)</a:t>",
                        slide_content,
                        re.DOTALL,
                    )
                    for match in title_matches:
                        try:
                            text = match.decode("utf-8", errors="replace").strip()
                            if text and not text.startswith("Click"):
                                paragraphs.append(text)
                        except Exception:
                            pass

                    # 2. Extract text from actual content placeholders
                    body_matches = re.findall(
                        b"<p:bodyPr[^>]*>.*?<a:t[^>]*>(.*?)</a:t>",
                        slide_content,
                        re.DOTALL,
                    )
                    for match in body_matches:
                        try:
                            text = match.decode("utf-8", errors="replace").strip()
                            if text and not text.startswith("Click"):
                                paragraphs.append(text)
                        except Exception:
                            pass

                    # If we found content for this slide, add it with the slide number
                    if paragraphs:
                        slide_text = f"Slide {i}:\n" + "\n".join(paragraphs)
                        slide_texts.append(slide_text)
            except Exception:
                continue

        # Return organized slide content
        if slide_texts:
            return "\n\n".join(slide_texts)

    finally:
        if close_zip:
            zip_file.close()

    return "[Could not extract text from PowerPoint presentation]"


def extract_text_from_xlsx(binary_data: bytes) -> str:
    try:
        excel_file = io.BytesIO(binary_data)
        xls = pd.ExcelFile(excel_file)
        csv_data = []
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            csv_text = df.to_csv(index=False, sep="|")
            csv_data.append(f"### {sheet_name} ###\n{csv_text}")
        return "Excel sheet data: " + "\n\n".join(csv_data)
    except Exception as e:
        logger.error(f"Error extracting data from Excel: {e}")
        return "[Could not extract text from Excel spreadsheet]"


def extract_text_from_pdf(binary_data: bytes) -> str:
    """
    Attempt to extract text from a PDF file.
    Falls back to basic text extraction if PDF-specific extraction fails.

    Args:
        binary_data: The binary data of the PDF file

    Returns:
        Extracted text from the PDF
    """
    # Check if the file looks like a PDF (starts with %PDF)
    if binary_data.startswith(b"%PDF"):
        try:
            # Try to extract text using a simple method that doesn't require additional libraries
            # This extracts all strings from the PDF that are likely to be text content
            result = []

            # Look for text objects in the PDF
            # PDF text is often in the format: BT ... (text) ... ET
            text_chunks = re.findall(rb"BT.*?ET", binary_data, re.DOTALL)

            for chunk in text_chunks:
                # Extract text strings in parentheses - these are usually content
                strings = re.findall(rb"\((.*?)\)", chunk, re.DOTALL)
                for s in strings:
                    try:
                        # Decode PDF-encoded strings
                        decoded = (
                            s.replace(b"\\(", b"(")
                            .replace(b"\\)", b")")
                            .decode("latin-1", errors="replace")
                        )
                        # Only keep strings that seem like real text (not just single characters)
                        if len(decoded.strip()) > 3 and not decoded.strip().isdigit():
                            result.append(decoded.strip())
                    except Exception:
                        pass

            # If we found text in the PDF, return it
            if result:
                return "\n".join(result)

        except Exception as e:
            logger.warning(f"Error extracting text from PDF using direct method: {e}")

    # Fall back to regular text extraction if PDF-specific extraction fails
    return binary_to_text(binary_data)


def file_extension_to_type(filename: str) -> str:
    """
    Determine file type based on file extension.

    Args:
        filename: Name of the file including extension

    Returns:
        String representing the type of the file
    """
    name = filename.lower()
    if name.endswith((".docx", ".doc")):
        return "document"
    elif name.endswith((".xlsx", ".xls")):
        return "spreadsheet"
    elif name.endswith((".pptx", ".ppt")):
        return "presentation"
    elif name.endswith((".pdf")):
        return "pdf"
    elif name.endswith((".txt", ".md", ".html", ".htm")):
        return "text"
    elif name.endswith((".jpg", ".jpeg", ".png", ".gif", ".bmp")):
        return "image"
    else:
        return "file"


def detect_content_type(content_type: str, filename: str) -> str:
    """
    Detect the appropriate content type based on content_type header and filename.

    Args:
        content_type: The content type header from HTTP response
        filename: Name of the file

    Returns:
        Normalized content type string
    """
    content_type = content_type.lower()

    # Direct matches for common types
    if "application/json" in content_type:
        return "json"
    elif "text/" in content_type:
        return "text"
    elif "application/xml" in content_type or content_type.endswith("+xml"):
        return "xml"
    elif "application/pdf" in content_type:
        return "pdf"
    elif "image/" in content_type:
        return "image"

    # Office document types
    if "application/vnd.openxmlformats-officedocument.wordprocessingml" in content_type:
        return "docx"
    elif "application/vnd.openxmlformats-officedocument.spreadsheetml" in content_type:
        return "xlsx"
    elif "application/vnd.openxmlformats-officedocument.presentationml" in content_type:
        return "pptx"
    elif "application/msword" in content_type:
        return "doc"
    elif "application/vnd.ms-excel" in content_type:
        return "xls"
    elif "application/vnd.ms-powerpoint" in content_type:
        return "ppt"

    # Fallback to file extension
    return file_extension_to_type(filename)


def process_sharepoint_response(
    response_content: bytes, content_type: str, filename: str
) -> Tuple[str, str]:
    """
    Process a SharePoint API response based on content type.

    Args:
        response_content: Binary content from response
        content_type: Content type header from the response
        filename: Name of the file

    Returns:
        Tuple of (extracted text, detected content type)
    """
    detected_type = detect_content_type(content_type, filename)

    if detected_type == "json":
        try:
            import json

            data = json.loads(response_content.decode("utf-8"))
            return json.dumps(data, indent=2), "application/json"
        except Exception as e:
            logger.warning(f"Failed to parse JSON: {e}")
            return binary_to_text(response_content), content_type

    elif detected_type in ["text", "xml"]:
        try:
            return response_content.decode("utf-8", errors="replace"), content_type
        except Exception:
            return binary_to_text(response_content), content_type

    elif detected_type == "pdf":
        return extract_text_from_pdf(response_content), content_type

    elif detected_type in [
        "docx",
        "xlsx",
        "pptx",
        "doc",
        "xls",
        "ppt",
        "document",
        "spreadsheet",
        "presentation",
    ]:
        # Handle Office documents using the existing functions
        if detected_type in ["docx", "document"]:
            return extract_text_from_docx(response_content), content_type
        elif detected_type in ["xlsx", "spreadsheet"]:
            return extract_text_from_xlsx(response_content), content_type
        elif detected_type in ["pptx", "presentation"]:
            return extract_text_from_pptx(response_content), content_type
        else:
            return binary_to_text(response_content), content_type

    else:
        # For any other file type, attempt to extract text
        return binary_to_text(response_content), content_type


def get_metadata_from_sharepoint_item(item: Dict) -> Dict:
    """
    Extract useful metadata from a SharePoint item.

    Args:
        item: SharePoint item data

    Returns:
        Dictionary of metadata
    """
    metadata = {
        "id": item.get("id", ""),
        "name": item.get("name", ""),
        "type": "folder"
        if item.get("folder")
        else file_extension_to_type(item.get("name", "")),
        "web_url": item.get("webUrl", ""),
        "created_by": item.get("createdBy", {}).get("user", {}).get("displayName", ""),
        "created_at": item.get("createdDateTime", ""),
        "modified_by": item.get("lastModifiedBy", {})
        .get("user", {})
        .get("displayName", ""),
        "modified_at": item.get("lastModifiedDateTime", ""),
    }

    # Add folder-specific metadata
    if item.get("folder"):
        metadata["child_count"] = item.get("folder", {}).get("childCount", 0)

    # Add file-specific metadata
    if item.get("file"):
        metadata["size"] = item.get("size", 0)
        metadata["mime_type"] = item.get("file", {}).get("mimeType", "")

    return metadata
