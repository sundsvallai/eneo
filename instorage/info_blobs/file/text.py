import asyncio
from enum import Enum
from pathlib import Path

import magic
import pptx
from docx2python import docx2python
from pypdf import PdfReader

from instorage.main.exceptions import FileNotSupportedException


class TextMimeTypes(str, Enum):
    MD = "text/markdown"
    TXT = "text/plain"
    PDF = "application/pdf"
    DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    TEXT_CSV = "text/csv"
    APP_CSV = "application/csv"
    PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    @classmethod
    def has_value(cls, value) -> bool:
        return any(value == item.value for item in cls)

    @classmethod
    def values(cls):
        return [item.value for item in cls]


AUDIO_FILE_TYPES = [
    "audio/x-m4a",
    "audio/ogg",
    "audio/wav",
    "audio/mpeg",
    "audio/mp3",
]


def _to_async(func):
    async def _inner(filepath: Path):
        return await asyncio.to_thread(func, filepath)

    return _inner


class TextExtractor:
    @staticmethod
    @_to_async
    def extract_from_plain_text(filepath: Path) -> str:
        return filepath.read_text("utf-8")

    @staticmethod
    @_to_async
    def extract_from_pdf(filepath: Path) -> str:
        reader = PdfReader(filepath)
        return " ".join([page.extract_text() for page in reader.pages])

    @staticmethod
    @_to_async
    def extract_from_docx(filepath: Path) -> str:
        with docx2python(filepath) as docx_content:
            return docx_content.text

    @staticmethod
    @_to_async
    async def extract_from_pptx(filepath: Path) -> str:
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(filepath)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"

    async def extract(self, filepath: Path):
        with open(filepath, "rb") as f:
            mimetype = magic.from_buffer(f.read(2048), mime=True)

        match mimetype:
            case (
                TextMimeTypes.TXT
                | TextMimeTypes.MD
                | TextMimeTypes.TEXT_CSV
                | TextMimeTypes.APP_CSV
            ):
                return await self.extract_from_plain_text(filepath)
            case TextMimeTypes.PDF:
                return await self.extract_from_pdf(filepath)
            case TextMimeTypes.DOCX:
                return await self.extract_from_docx(filepath)
            case TextMimeTypes.PPTX:
                return await self.extract_from_pptx(filepath)
            case _:
                raise FileNotSupportedException(f"{mimetype} not supported.")
